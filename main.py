import os
import json
from pathlib import Path
from datetime import datetime
import google.generativeai as genai
from pptx import Presentation
from dotenv import load_dotenv

# 載入環境變數 (例如 Gemini API Key)
load_dotenv()

class Code2Course:
    def __init__(self, target_dir=None, config=None):
        self.target_dir = target_dir
        self.config = config or {}
        self.api_key = self.config.get("GEMINI_API_KEY", "").strip()
        
        if not self.api_key:
            self.api_key = os.getenv("GEMINI_API_KEY", "").strip()
            
        if not self.api_key:
            raise ValueError("請確認 code2course_config.json 或 .env 檔案中已設定 GEMINI_API_KEY")
        
        # 設定 Gemini API Key
        genai.configure(api_key=self.api_key)
        # 初始化 Gemini 模型
        self.model = genai.GenerativeModel('gemini-2.5-flash')

    def _read_source_files(self):
        """讀取目標資料夾內的特定副檔名檔案"""
        if not self.target_dir:
            return {}
            
        supported_exts = self.config.get("supported_exts", ['.md', '.c', '.py', '.cpp', '.h'])
        contents = {}
        target_path = Path(self.target_dir)
        
        # 排除清單：避免讀取龐大的底層函式庫，防止超出 Gemini 100 萬 Token 限制
        exclude_dirs = self.config.get("exclude_dirs", ['drivers', 'rte', 'cmsis', 'build', 'node_modules', 'lib'])
        
        total_chars = 0
        max_chars = 300000  # 安全保護機制：限制最多讀取 30 萬字元，防止方案超用

        # 遞迴尋找符合的檔案
        if os.path.isdir(target_path):
            for ext in supported_exts:
                for filepath in target_path.rglob(f"*{ext}"):
                    file_str = str(filepath)
                    
                    if any(f"/{ex_dir}/" in file_str.lower() or f"\\{ex_dir}\\" in file_str.lower() for ex_dir in exclude_dirs):
                        continue
                        
                    if total_chars > max_chars:
                        print(f"\n⚠️ 目錄規模龐大，讀取已達到 Token 安全邊界 ({max_chars} 字元)。已自動略過剩餘的細節檔案。\n")
                        return contents

                    try:
                        # 加入 errors='ignore' 自動略過亂碼與舊編碼
                        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                            content = f.read()
                            if len(content) > 50000:
                                content = content[:50000] + "\n...[內容過長已截斷]"
                            contents[file_str] = content
                            total_chars += len(content)
                    except Exception as e:
                        pass
        return contents

    def stage1_scan_project(self, temp_dir):
        """
        第一階段：單純掃描專案並產生 Prompt (不呼叫 Gemini)
        功能：
        1. 走訪指定的 local 資料夾 (self.target_dir)
        2. 將程式碼與外層的 auto_edu_prompt.txt 結合
        3. 將生成的最終 Prompt 寫入到暫存區
        """
        print(f"-> [階段 1] 開始掃描資料夾: {self.target_dir}")
        source_texts = self._read_source_files()
        
        if not source_texts:
            print("未找到任何支援的程式碼或文件。")
            return None

        print("開始組合檔案內容...")
        merged_text = ""
        for filepath, content in source_texts.items():
            file_name = os.path.basename(filepath)
            merged_text += f"\n\n--- 檔案：{file_name} ---\n{content}\n"
            
        prompt_template_path = os.path.join(self.config.get("root_dir", "."), "code2course_prompt.txt")
        
        default_prompt = """[請注意：如果你要自訂此樣板，請保留最後一行的 {merged_text}，自動化工具才能把專案程式碼合併進來]

請你扮演「首席大師級技術講師 🧑‍🏫」，你不只是把冷冰冰的程式碼貼上去，而是要負責將這些包含多個章節、模組的專案原始碼，整合成一堂「有頭有尾、脈絡連貫」的線上微課程教材。

這份專案內含多個可能看似不相關的技術區塊，請你務必做到以下這三大要求：
1. **全局開場 (Course Overview)**：在課程最開頭，必須有一張投影片做「課程總覽與學習目標」，向學生說明這整套專案的最終目的，把各個模組串聯成一個有意義的大故事。
2. **段落連貫 (Transitions)**：在進入每個子模組時，請在講稿中發揮想像力，說明「剛剛我們學了什麼元件或概念，為什麼現在要接著學這個？」，讓各個技術區塊之間有嚴謹的邏輯關聯，絕不能是死板的列點。
3. **完美總結 (Conclusion)**：在課程最後，必須有一張投影片做「成果大總結」，回顧今天大家用這些模組共同拼湊出了什麼強大的架構或系統。

進入每個觀念的教學細節時，請在心中默默遵循這四個邏輯主軸來分析（但請根據專案實際內容替換掉這些起承轉合的名稱，絕對不可直接在標題印出 Scope 等字眼）：
- (Scope) 這是什麼核心功能？
- (Why do) 為何需要它？解決了什麼痛點或是工程難題？
- (How do) 運作邏輯是什麼？(善用生動的比喻或 Emoji 帶出畫面感)
- (What do) 程式碼具體達成了什麼？

輸出的格式必須是純 JSON 陣列 (Array) 結構，不需要 markdown code block (例如 ```json)，直接輸出 JSON 字串即可：
[
    {
        "slide_num": 1,
        "title": "🌟 課程總覽：探索強大的軟體架構",
        "bullet_points": ["從基礎建設到進階功能的完整旅程", "一次掌握多模組的強大火力"],
        "script": "各位同學們大家好！今天這堂課呢非常豐富，我們會看到許多不同的模組。大家可能會想，這些模組湊在一起能做什麼？其實，這就像是在打造產品的骨幹一樣..."
    },
    {
        "slide_num": 2,
        "title": "🤔 核心元件：如何讓系統運轉起來？",
        "bullet_points": ["系統開發中的關鍵樞紐", "用簡單的比喻看懂複雜邏輯"],
        "script": "回顧剛剛提到的全局概念，既然我們有了大方向，程式要怎麼具體執行呢？其實這時候就需要..."
    },
    {
        "slide_num": 3,
        "title": "🎓 總結回顧：見證系統的誕生",
        "bullet_points": ["回顧今天整合的各項底層技術", "未來延伸的無限應用可能"],
        "script": "課程到了尾聲，今天我們成功把剛才學的演算法與所有邏輯完美結合在一起了！恭喜大家完成這個專案的初步探索..."
    }
]

以下為專案原始資料內容：
{merged_text}
"""
        # 如果使用者已定義 auto_edu_prompt.txt 則讀取，否則自動產生預設模板
        if os.path.exists(prompt_template_path):
            try:
                with open(prompt_template_path, "r", encoding="utf-8") as pf:
                    prompt = pf.read()
            except Exception:
                prompt = default_prompt
        else:
            prompt = default_prompt
            try:
                with open(prompt_template_path, "w", encoding="utf-8") as pf:
                    pf.write(default_prompt)
                print(f"-> 初次執行，已生成樣板 Prompt 檔案: {prompt_template_path}")
            except Exception:
                pass
                
        # 安全機制：若使用者自訂的 Prompt 忘記加上 {merged_text}，則強力附加在最後
        if "{merged_text}" not in prompt:
            prompt += "\n\n以下為專案原始資料內容：\n{merged_text}"
            
        # 替換進原本的合併程式碼
        prompt = prompt.replace("{merged_text}", merged_text)
        
        # 輸出至紀錄檔案
        if temp_dir:
            try:
                prompt_path = os.path.join(temp_dir, "gemini_prompt_record.txt")
                with open(prompt_path, "w", encoding="utf-8") as f:
                    f.write(prompt)
                print(f"-> [階段 1] 已掃描完成！並將發送給 API 的 Prompt 封裝備份至: {prompt_path}")
            except Exception as e:
                print(f"無法寫入 Prompt：{e}")
                
        return prompt

    def stage2_gemini_and_ppt(self, prompt_text, final_dir, temp_dir, project_name):
        """
        第二階段：讀取已有的 prompt 字串，呼叫 Gemini 產出 JSON 分鏡圖與 PPT。
        """
        print("-> [階段 2] 發送請求至 Gemini API 中，請稍候...(分析大量內容可能需要數十秒)")
        try:
            response = self.model.generate_content(
                prompt_text,
                generation_config={"response_mime_type": "application/json"}
            )
            structured_content = json.loads(response.text)
            print(f"Gemini API 成功回傳！共產出 {len(structured_content)} 張投影片腳本。")
            
            # 將中間產生的結構化分析結果，另外存一份 JSON 在暫存區給使用者備查
            if structured_content and temp_dir:
                analysis_file = os.path.join(temp_dir, "gemini_analysis_result.json")
                with open(analysis_file, "w", encoding="utf-8") as f:
                    json.dump(structured_content, f, ensure_ascii=False, indent=4)
                print(f"-> 分析腳本分鏡資料保存至: {analysis_file}")
                
            # Stage 2: 產生 PPT 簡報
            if final_dir:
                ppt_file = os.path.join(final_dir, f"{project_name}_course_slides.pptx")
                self._generate_internal_ppt(structured_content, output_ppt_path=ppt_file)
                print(f"-> [階段 2] PPT 簡報檔成功保存至: {ppt_file}")
            else:
                ppt_file = None
                
            return structured_content, ppt_file
        except json.JSONDecodeError as decode_err:
            print(f"解析 JSON 發生錯誤，回傳原始文字可能非標準 JSON 格式：{decode_err}")
            return None, None
        except Exception as e:
            print(f"呼叫 Gemini API 發生錯誤: {e}")
            return None, None

    def _generate_internal_ppt(self, structured_content, output_ppt_path):
        if not structured_content:
            return
        prs = Presentation()
        for slide_data in structured_content:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            body_shape = slide.shapes.placeholders[1]
            
            title_shape.text = str(slide_data.get('title', '無標題'))
            tf = body_shape.text_frame
            for point in slide_data.get('bullet_points', []):
                p = tf.add_paragraph()
                p.text = str(point)

        prs.save(output_ppt_path)

    def stage3_synthesize_media(self, structured_content, output_video_path, temp_dir):
        """
        第三階段：語音與影片合成
        功能：
        1. 語音：使用 edge-tts (subprocess) 將腳本轉換為語音檔
        2. 影片：從結構化分鏡表繪製背景 (使用 Pillow)，與語音檔合成 (MoviePy)
        """
        try:
            # 兼容 MoviePy 1.x 代碼風格
            from moviepy.editor import ImageClip, AudioFileClip, concatenate_videoclips
        except ImportError:
            # 相容 MoviePy 2.x 代碼風格
            from moviepy import ImageClip, AudioFileClip, concatenate_videoclips
            
        import subprocess
        from PIL import Image, ImageDraw, ImageFont

        if not structured_content:
            print("無結構化 JSON 腳本可進行影片合成。")
            return

        print("-> [階段 3] 開始進行語音合成與教學影片生成...")
        
        final_clips = []
        for slide_data in structured_content:
            slide_num = slide_data.get('slide_num', 1)
            text_script = str(slide_data.get('script', ''))
            audio_path = os.path.join(temp_dir, f"slide_{slide_num}_audio.mp3")
            img_path = os.path.join(temp_dir, f"slide_{slide_num}_bg.png")
            
            print(f"正在生成第 {slide_num} 頁的語音 (edge-tts)...")
            # 使用 edge-tts CLI 工具生成音檔
            subprocess.run(["edge-tts", "--voice", "zh-TW-HsiaoChenNeural", "--text", text_script, "--write-media", audio_path], check=True)
            
            # 使用 Pillow 繪製靜態背景圖
            img = Image.new('RGB', (1920, 1080), color=(25, 28, 36))
            draw = ImageDraw.Draw(img)
            
            # 裝飾性頂部橫幅
            draw.rectangle([(0, 0), (1920, 150)], fill=(41, 128, 185))
            
            # 尋找 Mac 與 Linux (Ubuntu CI) 常見支援中文的字型
            font_title = None
            font_body = None
            system_fonts = [
                # macOS fonts
                "/System/Library/Fonts/PingFang.ttc",
                "/System/Library/Fonts/STHeiti Light.ttc",
                "/Library/Fonts/Arial Unicode.ttf",
                "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
                # Ubuntu / GitHub Actions fonts
                "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
                "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
                "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc"
            ]
            
            for f_path in system_fonts:
                if os.path.exists(f_path):
                    try:
                        font_title = ImageFont.truetype(f_path, 80)
                        font_body = ImageFont.truetype(f_path, 50)
                        break
                    except Exception:
                        continue
                        
            if font_title is None or font_body is None:
                print("⚠️ 無法載入中文支援字型，影片文字可能會變成方塊亂碼。")
                font_title = ImageFont.load_default()
                font_body = ImageFont.load_default()
            
            # 繪製標題 (對齊橫幅)
            title_text = str(slide_data.get('title', '無標題'))
            draw.text((100, 35), title_text, font=font_title, fill=(255, 255, 255))
            
            # 繪製重點 (支援簡單換行與項目符號)
            y_offset = 250
            for point in slide_data.get('bullet_points', []):
                bullet_str = f"- {point}"
                chunk_len = 35
                chunks = [bullet_str[i:i+chunk_len] for i in range(0, len(bullet_str), chunk_len)]
                for chunk in chunks:
                    draw.text((120, y_offset), chunk, font=font_body, fill=(210, 220, 240))
                    y_offset += 70
                y_offset += 30
                
            img.save(img_path)
            
            # 結合 MoviePy
            audioclip = AudioFileClip(audio_path)
            videoclip = ImageClip(img_path)
            
            # 動態兼容 MoviePy 1.x 與 2.x 的 API
            if hasattr(videoclip, 'with_duration'):
                videoclip = videoclip.with_duration(audioclip.duration)
                videoclip = videoclip.with_audio(audioclip)
            else:
                videoclip = videoclip.set_duration(audioclip.duration)
                videoclip = videoclip.set_audio(audioclip)
                
            final_clips.append(videoclip)
            
        print("所有語音與分頁畫面生成完畢，開始合成最終影片...")
        final_video = concatenate_videoclips(final_clips)
        final_video.write_videofile(output_video_path, fps=24, codec="libx264", audio_codec="aac")
        
        print(f"-> [階段 3] 教學影片合成完畢，已保存至: {output_video_path}")


# =========================================================
# 選單與流程控制器 (CLI Runner)
# =========================================================

def _get_target_folders(root_dir):
    try:
        folders = [f for f in os.listdir(root_dir) 
                   if os.path.isdir(os.path.join(root_dir, f)) and not f.startswith('.')]
        folders.sort()
        return folders
    except Exception:
        return []

def select_history_stage_dir(root_dir):
    """用於階段 2 與 3，從 outputs 選出前人留下的 scope 與 time 資料夾進度"""
    import questionary
    outputs_dir = os.path.join(root_dir, "outputs")
    if not os.path.exists(outputs_dir):
        print("📁 找不到 outputs 資料夾，請先執行 階段 1。")
        return None, None
        
    projects = [f for f in os.listdir(outputs_dir) if os.path.isdir(os.path.join(outputs_dir, f))]
    if not projects:
        print(" outputs 內無任何專案歷史，請先執行階段 1。")
        return None, None
    
    selected_project = questionary.select(
        "📂 請選擇要接續處理的專案 (scope_folder)：", 
        choices=projects
    ).ask()
    if not selected_project:
        return None, None
    
    project_path = os.path.join(outputs_dir, selected_project)
    timestamps = [f for f in os.listdir(project_path) if os.path.isdir(os.path.join(project_path, f))]
    timestamps.sort(reverse=True) # 近期的放前面
    if not timestamps:
        print(f"專案 {selected_project} 找不到任何歷史時間點資料夾。")
        return None, None
    
    selected_timestamp = questionary.select(
        "🕒 請選擇要接續執行進度的時間點 (time folder)：", 
        choices=timestamps
    ).ask()
    if not selected_timestamp:
        return None, None
    
    return selected_project, selected_timestamp

def main():
    print("====== 🤖 Code2Course 教學影片生成工具 (多模式版) ======")
    
    current_dir = os.path.dirname(os.path.abspath(__file__))
    root_dir = os.path.dirname(current_dir)
    config_path = os.path.join(root_dir, 'code2course_config.json')
    
    # 嘗試讀取 external config
    config = {}
    if os.path.exists(config_path):
        try:
            with open(config_path, "r", encoding="utf-8") as f:
                config = json.load(f)
        except Exception as e:
            print(f"讀取設定檔發生錯誤: {e}")
    config['root_dir'] = root_dir

    try:
        import questionary
    except ImportError:
        print("\n缺少 questionary 套件，請先在終端機輸入: pip install -r requirements.txt\n")
        return

    # 要求使用者選擇模式
    mode_choices = [
        questionary.Choice(title="[0] 完整自動化流程 (掃描 -> Gemini 分析 -> 產出影片)", value=0),
        questionary.Choice(title="[1] 獨立階段 1：純粹掃描專案，產生 Prompt 檔 (供您手動微調)", value=1),
        questionary.Choice(title="[2] 獨立階段 2：讀出已處理的 Prompt (接關)，送交 Gemini 產生 PPT 與 Json 腳本", value=2),
        questionary.Choice(title="[3] 獨立階段 3：讀出已產生的 Json 腳本 (接關)，進行 TTS 合成語音與背景影片", value=3)
    ]
    
    # 支援 Github Actions 的無介面自動執行 (直接吃環境變數)
    env_mode = os.getenv("AUTO_MODE")
    if env_mode is not None and env_mode.strip().isdigit():
        selected_mode = int(env_mode.strip())
        print(f"🤖 偵測到 CI 變數 AUTO_MODE={selected_mode}，跳過選單直接執行。")
    else:
        selected_mode_obj = questionary.select(
            "請選擇要執行的功能模組：",
            choices=mode_choices
        ).ask()
        
        if selected_mode_obj is None:
            print("已取消選擇。")
            return
        selected_mode = selected_mode_obj

    target_directory = None
    project_name = None
    timestamp_str = None
    
    # 如果是模式 0 或是 1 (需要掃描)，那麼才需要去挑選外部專案資料夾
    if selected_mode in [0, 1]:
        target_folder = config.get("target_folder", "").strip()

        if target_folder:
            # Setting 檔強制覆寫
            if target_folder == os.path.basename(root_dir) or target_folder == ".":
                target_directory = root_dir
                project_name = os.path.basename(root_dir)
            else:
                target_directory = os.path.join(root_dir, target_folder)
                project_name = target_folder
            print(f"✅ Config 設定檔強制指定掃描目標: {project_name}")
        else:
            folders = _get_target_folders(root_dir)
            if not folders:
                print(f"在 {root_dir} 找不到任何可選資料夾。")
                return

            choices = [f"📁 {f}" for f in folders]
            choices.insert(0, f"📂 (掃描整個根目錄) {os.path.basename(root_dir)}")
            
            selected = questionary.select(
                "請使用上下鍵選擇要掃描的專案資料夾：",
                choices=choices
            ).ask()
            
            if not selected:
                return
                
            if selected.startswith("📂"):
                target_directory = root_dir
                project_name = os.path.basename(root_dir)
            else:
                folder_name = selected.replace("📁 ", "")
                target_directory = os.path.join(root_dir, folder_name)
                project_name = folder_name
                
        # 新建立一層「本次執行的時間戳」，以便未來接關使用
        timestamp_str = datetime.now().strftime("%Y%m%d_%H%M%S")

    else:
        # 如果是模式 2 或 3，程式將由 outputs 內取得接續的 project_name 與 timestamp_str
        proj, ts = select_history_stage_dir(root_dir)
        if not proj or not ts:
            return
        project_name = proj
        timestamp_str = ts

    # ======= 路徑組裝 =======
    output_base_dir = os.path.join(root_dir, "outputs", project_name, timestamp_str)
    final_output_dir = os.path.join(output_base_dir, "final")
    temp_intermediate_dir = os.path.join(output_base_dir, "temp_analysis_and_tts")
    
    prompt_record_path = os.path.join(temp_intermediate_dir, "gemini_prompt_record.txt")
    json_analysis_path = os.path.join(temp_intermediate_dir, "gemini_analysis_result.json")

    # 如果是模式 0 或 1 (產生新目錄)，要確實建立實體目錄
    if selected_mode in [0, 1]:
        os.makedirs(final_output_dir, exist_ok=True)
        os.makedirs(temp_intermediate_dir, exist_ok=True)
        print(f"\n✅ 準備掃描目標資料夾: {target_directory}")
        if selected_mode == 0:
            print(f"📁 最終產出將儲存於: outputs/{project_name}/{timestamp_str}/final/")
        print(f"📁 中間暫存檔將儲存於: outputs/{project_name}/{timestamp_str}/temp_analysis_and_tts/\n")
    else:
        # 接龍模式，驗證目錄確實存在
        if not os.path.exists(temp_intermediate_dir):
            print(f"❌ 查無此進度的暫存資料夾，請確認進度正確。({temp_intermediate_dir})")
            return

    # ======= 宣告工具實體 =======
    try:
        # 對於階段 2, 3 若 target_directory 為 None 亦無妨，掃描檔不會被調用到
        generator = Code2Course(target_dir=target_directory, config=config)
    except ValueError as e:
        print(f"👉 啟動提醒: {e}")
        return

    # ======= 開始執行邏輯流程 =======

    # ▶️ [階段 1] 啟動
    if selected_mode == 0 or selected_mode == 1:
        prompt_text = generator.stage1_scan_project(temp_dir=temp_intermediate_dir)
        
        if selected_mode == 1:
            print("\n=== ✨ [階段 1] 單獨執行結束！ ===")
            print(f"請前往修改: \n{prompt_record_path}\n確認 Prompt 無誤後，隨時可啟動模式 [2] 進入接關。")
            return

    # ▶️ [階段 2] 啟動
    if selected_mode == 0 or selected_mode == 2:
        if selected_mode == 2:
            print(f"\n✅ [階段 2] 接關：從紀錄檔讀取 Prompt\n檔案: {prompt_record_path}")
            if not os.path.exists(prompt_record_path):
                print("❌ 找不到 Prompt 紀錄檔。您必須先針對此專案執行階段 1！")
                return
            with open(prompt_record_path, "r", encoding="utf-8") as f:
                prompt_text = f.read()

        structured_content, ppt_file = generator.stage2_gemini_and_ppt(
            prompt_text=prompt_text,
            final_dir=final_output_dir,
            temp_dir=temp_intermediate_dir,
            project_name=project_name
        )
        
        if selected_mode == 2:
            print("\n=== ✨ [階段 2] 單獨執行結束！ ===")
            if structured_content:
                print(f"請確認 PPT 結構（已存放至 {final_output_dir} / 或可查看產生出的 json 腳本）\n確認無誤後，可啟動模式 [3] 開始生成影音。")
            return

    # ▶️ [階段 3] 啟動
    if selected_mode == 0 or selected_mode == 3:
        if selected_mode == 3:
            print(f"\n✅ [階段 3] 接關：讀取 JSON 分鏡腳本，進入影音合成")
            if not os.path.exists(json_analysis_path):
                print(f"❌ 找不到腳本 JSON 檔: {json_analysis_path}。請先執行階段 2 產生分鏡。")
                return
            with open(json_analysis_path, "r", encoding="utf-8") as f:
                structured_content = json.load(f)

        video_file = os.path.join(final_output_dir, f"{project_name}_course_video.mp4")
        generator.stage3_synthesize_media(
            structured_content=structured_content, 
            output_video_path=video_file,
            temp_dir=temp_intermediate_dir
        )
        
        if selected_mode == 3:
            print("\n=== ✨ [階段 3] 單獨執行結束！ ===")
            print(f"完整影片與語音素材已存放至 {final_output_dir} 內。")
            return

    # 完整模式結束
    if selected_mode == 0:
        print("\n=== 🎯 [所有流程自動化執行完成] 🎯 ===")

if __name__ == "__main__":
    main()
