import os
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement

def add_text_to_slide(slide, slide_data, style, color, s_width, s_height, centered=False, width=None, show_body_mask=False):
    """輔助函式：將標題與內文加入投影片"""
    text_width = width or s_width - Inches(1.0)
    left_margin = (s_width - text_width) / 2 if centered else Inches(0.5)
    
    # 標題
    title_text = slide_data.get('title', '無標題')
    title_top = Inches(style.get("title_top_margin_inches", 0.5))
    title_box = slide.shapes.add_textbox(left_margin, title_top, text_width, Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True # 允許標題換行，避免超出畫面
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.bold = True
    
    # 動態調整標題字級與預估行數 (從 style 設定讀取字級)
    estimated_lines = 1
    if len(title_text) > 35:
        p.font.size = Pt(style.get("title_font_size_small", 28))
        estimated_lines = 2 if len(title_text) < 70 else 3
    elif len(title_text) > 18:
        p.font.size = Pt(style.get("title_font_size_medium", 36))
        estimated_lines = 2
    else:
        p.font.size = Pt(style.get("title_font_size_large", 44))
        estimated_lines = 1
        
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT

    # 動態計算內文起始位置 (body_top)，給予更多緩衝空間 (0.85" per line)
    dynamic_body_top = title_top + Inches(estimated_lines * 0.85)
    body_top = max(dynamic_body_top, Inches(style.get("body_top_margin_min_inches", 2.5)))
    body_height = s_height - body_top - Inches(0.5)
    
    # 如果需要針對內文加入遮罩
    if show_body_mask:
        # 加入半透明深灰色矩形 (從 style 讀取透明度與顏色)
        mask = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_margin - Inches(0.2), body_top - Inches(0.1), 
                                      text_width + Inches(0.4), body_height + Inches(0.2))
        mask.fill.solid()
        
        m_rgb = style.get("mask_rgb", [40, 40, 40])
        m_alpha = style.get("mask_alpha", 0.5)
        mask.fill.fore_color.rgb = RGBColor(*m_rgb)
        
        # 使用 XML 直接注入透明度
        try:
            alpha = OxmlElement('a:alpha')
            alpha.set('val', f'{int(m_alpha * 100000)}') 
            mask.fill._xPr.xpath('.//a:srgbClr')[0].append(alpha)
        except Exception:
            mask.fill.transparency = m_alpha # Fallback
            
        mask.line.fill.background() # 無邊框

    body_box = slide.shapes.add_textbox(left_margin, body_top, text_width, body_height)
    tf = body_box.text_frame
    tf.word_wrap = True
    
    for point in slide_data.get('bullet_points', []):
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(style.get("body_font_size", 24))
        p.font.color.rgb = color
        p.space_after = Pt(10)
        p.alignment = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT

def run_stage4(structured_content, output_ppt_path, style, temp_dir=None):
    """
    第四階段：產生/更新 PPT 簡報
    """
    print("-> [階段 4] 開始產生 PPT 簡報...")
    if not structured_content:
        return
        
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 讀取圖片來源資訊 (用於最後一頁)
    image_sources = {}
    if temp_dir:
        source_file = os.path.join(temp_dir, "image_sources.json")
        if os.path.exists(source_file):
            try:
                with open(source_file, "r", encoding="utf-8") as f:
                    image_sources = json.load(f)
            except Exception:
                pass

    for slide_data in structured_content:
        slide_num = slide_data.get('slide_num', 1)
        include_image = slide_data.get('include_image', True)
        
        # 檢查是否有實際下載的圖片
        pexels_img_path = os.path.join(temp_dir, f"slide_{slide_num}_pexels_bg.jpg") if temp_dir else None
        has_image = include_image and pexels_img_path and os.path.exists(pexels_img_path)

        # 使用空白佈局以便完全掌控元件位置
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        if has_image:
            # 統一採用全螢幕背景圖排版
            slide.shapes.add_picture(pexels_img_path, 0, 0, width=slide_width, height=slide_height)
            
            # 直接放置文字 (根據 style 設定決定顏色與是否加遮罩)
            txt_color_rgb = style.get("text_color_on_image_rgb", [255, 255, 255])
            add_text_to_slide(slide, slide_data, style, color=RGBColor(*txt_color_rgb), centered=False, 
                             s_width=slide_width, s_height=slide_height, show_body_mask=True)
        else:
            # 無背景圖的簡約白底排版 (顏色從 style 讀取)
            bg_color_rgb = style.get("bg_color_plain_rgb", [255, 255, 255])
            background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(*bg_color_rgb)
            background.line.fill.background()
            
            # 靠左文字 (顏色從 style 讀取)
            txt_color_rgb = style.get("text_color_on_plain_rgb", [0, 0, 0])
            add_text_to_slide(slide, slide_data, style, color=RGBColor(*txt_color_rgb), centered=False,
                             s_width=slide_width, s_height=slide_height)

    # 產生最後一頁：圖片來源標記
    if image_sources:
        ref_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # 深色背景
        bg = ref_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(20, 20, 20)
        bg.line.fill.background()

        # 標題
        title_box = ref_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), slide_width - Inches(1), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "🖼️ Image References & Attributions"
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.LEFT

        # 列表
        list_top = Inches(1.5)
        list_box = ref_slide.shapes.add_textbox(Inches(0.5), list_top, slide_width - Inches(1), slide_height - list_top - Inches(0.5))
        tf = list_box.text_frame
        tf.word_wrap = True

        for s_num, metadata in sorted(image_sources.items(), key=lambda x: int(x[0])):
            p = tf.add_paragraph()
            photographer = metadata.get('photographer', 'Unknown')
            p.text = f"- Slide {s_num}: Photo by {photographer} on Pexels"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(200, 200, 200)
        
        print(f"-> [階段 4] 已在 PPT 加入圖片來源標記頁 (共 {len(image_sources)} 筆資料)。")

    prs.save(output_ppt_path)
